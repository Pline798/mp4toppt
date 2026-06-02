"""文件生成器 - 将截图序列转换为 PPT 或 PDF"""

from typing import List
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF


class FileGenerator:
    """将图片列表生成为 PPTX 或 PDF 文件"""

    @staticmethod
    def _add_image_to_slide(slide, img_path: str, slide_width, slide_height):
        """将图片居中自适应添加到幻灯片"""
        img = Image.open(img_path)
        w, h = img.size
        aspect = w / h
        margin = Inches(0.5)
        max_w = slide_width - 2 * margin
        max_h = slide_height - 2 * margin

        if aspect > max_w / max_h:
            img_width = max_w
            img_height = max_w / aspect
        else:
            img_height = max_h
            img_width = max_h * aspect

        left = (slide_width - img_width) / 2
        top = (slide_height - img_height) / 2
        slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)

    @staticmethod
    def create_ppt(temp_images: List[str], output_path: str):
        """生成 PPTX 文件"""
        prs = Presentation()
        layout = prs.slide_layouts[5]  # blank
        sw, sh = prs.slide_width, prs.slide_height

        for img_path in temp_images:
            slide = prs.slides.add_slide(layout)
            FileGenerator._add_image_to_slide(slide, img_path, sw, sh)

        prs.save(output_path)

    @staticmethod
    def create_pdf(temp_images: List[str], output_path: str):
        """生成 PDF 文件"""
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        for img_path in temp_images:
            pdf.add_page()
            img = Image.open(img_path)
            w, h = img.size
            aspect = w / h
            pdf_w = pdf.w - 20
            pdf_h = pdf_w / aspect
            if pdf_h > pdf.h - 20:
                pdf_h = pdf.h - 20
                pdf_w = pdf_h * aspect
            x = (pdf.w - pdf_w) / 2
            y = (pdf.h - pdf_h) / 2
            pdf.image(img_path, x=x, y=y, w=pdf_w, h=pdf_h)

        pdf.output(output_path)

    @staticmethod
    def generate(temp_images: List[str], output_format: str, output_path: str):
        """统一入口"""
        if output_format == "pptx":
            FileGenerator.create_ppt(temp_images, output_path)
        else:
            FileGenerator.create_pdf(temp_images, output_path)
